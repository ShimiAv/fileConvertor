import os
import PyPDF2
import pandas
import pptx
import pymysql as sql
from docx import Document
from dotenv import load_dotenv
#fixed path to file

def main_process(keyword):
    load_dotenv("data.env")
    DB_HOST = os.getenv("DB_HOST")
    DB_USER = os.getenv("DB_USER")
    DB_PASSWORD = os.getenv("DB_PASSWORD")
    DB_NAME = os.getenv("DB_NAME")

    try:
        connection = sql.connect(
            host=DB_HOST,
            user=DB_USER,
            database=DB_NAME,
            password=DB_PASSWORD,
            port=3306,

        )
        cursor = connection.cursor()
        cursor.execute('SELECT id, url, name FROM files')
        records = cursor.fetchall()
        exists = []
        for record in records:
            print(record)
            if isExist(record[1] + "\\" + record[2], keyword):
                exists.append(record[0])
        return exists

    except sql.MySQLError as db_error:
        print(f"Database error occurred: {str(db_error)}")
    except Exception as e:
        print(f"Error occurred: {str(e)}")



def isExist(path, keyword):
    file_ending = path.split('.')[-1].lower()
    match file_ending:
        case 'csv':
            data = CSV_to_reader(path)
            return data is not None and keyword in str(data.values).lower()
        case 'docx':
            text = word_to_reader(path)
            return text and keyword in text.lower()
        case 'pptx':
            text = ppt_to_reader(path)
            return text and keyword in ' '.join(text).lower()
        case 'pdf':
            text = pdf_to_reader(path)
            return text and keyword in ' '.join(text).lower()
        case _:
            print(f"unsupported file format: {file_ending}")
            return False


def CSV_to_reader(path):
    try:
        data = pandas.read_csv(path)
        print(data)
        return data
    except Exception as e:
        print(f"An error occurred while reading the csv file : {e}")
        return None


def word_to_reader(path):
    try:
        doc = Document(path)
        text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        print(text_content)
        return text_content
    except Exception as e:
        print(f"An error occurred while reading the Word file: {e}")
        return None


def ppt_to_reader(path):
    presentation = pptx.Presentation(path)
    slide_texts = []

    try:
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_texts.append(run.text)
                        print(run.text)

        return slide_texts
    except Exception as e:
        print(f"An error occurred while reading the ppt file : {e}")
        return None


def pdf_to_reader(path):
    page_texts = []

    try:
        with open(path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                page_texts.append(page_text)
                print(page_text)

    except Exception as e:
        print(f"An error occurred while processing the PDF: {str(e)}")

    return page_texts

